import os
import json
import matplotlib.pyplot as plt
from typing import Any, List, Dict
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.shared import Inches as DocxInches
from langchain_core.tools import tool

OUTPUT_DIR = os.getenv("OUTPUT_DIR", os.path.join(os.getcwd(), "output"))
os.makedirs(OUTPUT_DIR, exist_ok=True)

@tool
def read_document(filename: str) -> str:
    """
    Reads the current text/data content of an existing DOCX, PPTX, or XLSX file.
    Always use this BEFORE modifying a document to know its current state and structure.
    """
    filepath = os.path.join(OUTPUT_DIR, filename)
    if not os.path.exists(filepath):
        return f"Error: File {filename} does not exist in {OUTPUT_DIR}."
    
    ext = filename.split('.')[-1].lower()
    try:
        if ext == 'docx':
            doc = Document(filepath)
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        elif ext == 'xlsx':
            import pandas as pd
            # Read all sheets into a dictionary of lists
            df_dict = pd.read_excel(filepath, sheet_name=None)
            result = {}
            for sheet, df in df_dict.items():
                result[sheet] = json.loads(df.to_json(orient="records"))
            return json.dumps(result)
        elif ext == 'pptx':
            prs = Presentation(filepath)
            content = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                content.append(f"Slide {i+1}: {' | '.join(slide_text)}")
            return "\n".join(content)
        else:
            return "Unsupported file type for reading."
    except Exception as e:
        return f"Error reading file: {str(e)}"

@tool
def generate_chart(title: str, labels: List[str], values: List[float], chart_type: str = "bar") -> str:
    """
    Generates a chart (bar or pie) and saves it as an image.
    Returns the file path of the generated chart.
    """
    plt.figure(figsize=(6, 4))
    if chart_type == "pie":
        plt.pie(values, labels=labels, autopct='%1.1f%%')
    else:
        plt.bar(labels, values)
    
    plt.title(title)
    plt.tight_layout()
    
    filename = f"{title.replace(' ', '_').lower()}.png"
    filepath = os.path.join(OUTPUT_DIR, filename)
    plt.savefig(filepath)
    plt.close()
    
    return filepath

@tool
def create_pptx(filename: str, slides_data: str) -> str:
    """
    Generates or overwrites a PPTX presentation.
    slides_data must be a JSON string representing a list of dictionaries:
    [{"title": "Slide 1", "bullets": ["Point 1", "Point 2"], "image_path": "optional/path.png"}]
    """
    try:
        data = json.loads(slides_data)
        prs = Presentation()
        
        for slide_info in data:
            slide_layout = prs.slide_layouts[1] 
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = slide_info.get("title", "")
            tf = body_shape.text_frame
            
            for bullet in slide_info.get("bullets", []):
                p = tf.add_paragraph()
                p.text = bullet
                
            if "image_path" in slide_info and os.path.exists(slide_info["image_path"]):
                slide.shapes.add_picture(slide_info["image_path"], Inches(5), Inches(2), width=Inches(4))
                
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)
        return f"Successfully created/updated PPTX at: {filepath}"
    except Exception as e:
        return f"Error creating PPTX: {str(e)}"

@tool
def create_docx(filename: str, title: str, sections_data: str) -> str:
    """
    Generates or overwrites a DOCX report.
    sections_data must be a JSON string representing a list of dictionaries:
    [{"heading": "Section 1", "content": "Paragraph text", "image_path": "optional/path.png"}]
    """
    try:
        data = json.loads(sections_data)
        doc = Document()
        doc.add_heading(title, 0)
        
        for section in data:
            doc.add_heading(section.get("heading", ""), level=1)
            doc.add_paragraph(section.get("content", ""))
            
            if "image_path" in section and os.path.exists(section["image_path"]):
                doc.add_picture(section["image_path"], width=DocxInches(5))
                
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return f"Successfully created/updated DOCX at: {filepath}"
    except Exception as e:
        return f"Error creating DOCX: {str(e)}"

@tool
def create_xlsx(filename: str, sheet_name: str, data: str) -> str:
    """
    Generates or overwrites an XLSX spreadsheet.
    data must be a JSON string representing a list of dictionaries.
    """
    try:
        import pandas as pd
        data_list = json.loads(data)
        df = pd.DataFrame(data_list)
        
        filepath = os.path.join(OUTPUT_DIR, filename)
        df.to_excel(filepath, index=False, sheet_name=sheet_name)
        return f"Successfully created/updated XLSX at: {filepath}"
    except Exception as e:
        return f"Error creating XLSX: {str(e)}"
    
class DocumentToolset:
    def get_tools(self) -> list:
        # Added read_document to the toolset
        return [read_document, generate_chart, create_pptx, create_docx, create_xlsx]